import PyPDF2
from pptx import Presentation
import re

class ContentProcess:
    def __init__(self):
        return
    
    def convertPDF(self, pdf_path:str)->str:
        extracted = ""
        with open(pdf_path, 'rb') as file:
            pdfR = PyPDF2.PdfReader(file)
            for page in range(len(pdfR.pages)):
                extracted += pdfR.pages[page].extract_text()
        return extracted
    
    def convertPPTX(self, pptx_path:str)->str:
        extracted = ""
        presentation = Presentation(pptx= pptx_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        extracted += para.text + '\n'
        return extracted
    
    def clean(self, extracted:str)->str:
        extracted = re.sub(r'\s+', ' ', extracted)
        extracted = re.sub(r'[^a-zA-Z0-9\s]', '', extracted)
        return extracted.strip()